import sys
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_deck(output_path):
    prs = Presentation()
    # 16:9 Widescreen dimensions
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_slide_layout = prs.slide_layouts[6]

    # Color Palette - Sanad Brand
    COLOR_BG_DARK = RGBColor(11, 11, 20)        # #0B0B14 Deep Dark
    COLOR_SURFACE = RGBColor(20, 20, 36)        # #141424 Card surface
    COLOR_SURFACE_ALT = RGBColor(28, 28, 50)    # #1C1C32 Lighter card
    COLOR_BRAND = RGBColor(107, 0, 255)         # #6B00FF Primary Purple
    COLOR_BRAND_LIGHT = RGBColor(155, 89, 255)  # #9B59FF Light Purple
    COLOR_BRAND_ACCENT = RGBColor(196, 153, 255)# #C499FF Lavender
    COLOR_TEXT_WHITE = RGBColor(255, 255, 255)
    COLOR_TEXT_MUTED = RGBColor(160, 160, 185)
    COLOR_GREEN = RGBColor(16, 185, 129)        # Emerald
    COLOR_AMBER = RGBColor(245, 158, 11)        # Amber
    COLOR_BLUE = RGBColor(59, 130, 246)         # Blue
    COLOR_RED = RGBColor(239, 68, 68)           # Red

    def set_slide_background(slide, color=COLOR_BG_DARK):
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = color
        bg.line.fill.background() # no line
        return bg

    def add_header(slide, tag_text, title_text, subtitle_text=""):
        # Tag pill
        tag_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(0.5), Inches(3.2), Inches(0.35))
        tag_box.fill.solid()
        tag_box.fill.fore_color.rgb = COLOR_SURFACE
        tag_box.line.color.rgb = COLOR_BRAND
        tag_box.line.width = Pt(1)
        tf = tag_box.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = tag_text
        p.font.name = "Arial"
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = COLOR_BRAND_ACCENT
        p.alignment = PP_ALIGN.CENTER

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.9), Inches(11.7), Inches(0.6))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.name = "Arial"
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = COLOR_TEXT_WHITE

        # Subtitle
        if subtitle_text:
            sub_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.45), Inches(11.7), Inches(0.4))
            tf = sub_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = subtitle_text
            p.font.name = "Arial"
            p.font.size = Pt(12)
            p.font.color.rgb = COLOR_TEXT_MUTED

    # ==========================================
    # SLIDE 1: COVER SLIDE
    # ==========================================
    slide1 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide1)

    # Ambient Card in Center
    center_card = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.2), Inches(1.0), Inches(10.9), Inches(5.5))
    center_card.fill.solid()
    center_card.fill.fore_color.rgb = COLOR_SURFACE
    center_card.line.color.rgb = COLOR_BRAND
    center_card.line.width = Pt(1.5)

    # Logo / Brand Pill
    logo_box = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.8), Inches(1.5), Inches(3.7), Inches(0.5))
    logo_box.fill.solid()
    logo_box.fill.fore_color.rgb = COLOR_BRAND
    logo_box.line.fill.background()
    p = logo_box.text_frame.paragraphs[0]
    p.text = "⚡ سند AI — SANAD AI"
    p.font.name = "Arial"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLOR_TEXT_WHITE
    p.alignment = PP_ALIGN.CENTER

    # Main Title
    t_box = slide1.shapes.add_textbox(Inches(1.5), Inches(2.3), Inches(10.3), Inches(1.2))
    tf = t_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "الخطة الاستراتيجية وتحليل السوق"
    p.font.name = "Arial"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COLOR_TEXT_WHITE
    p.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.text = "تحليل SWOT • دراسة المنافسين • أهداف الكوارتر القادم (OKRs)"
    p2.font.name = "Arial"
    p2.font.size = Pt(18)
    p2.font.bold = True
    p2.font.color.rgb = COLOR_BRAND_LIGHT
    p2.alignment = PP_ALIGN.CENTER

    # Meta Info Row
    meta_box = slide1.shapes.add_textbox(Inches(1.5), Inches(4.8), Inches(10.3), Inches(1.0))
    tf = meta_box.text_frame
    p = tf.paragraphs[0]
    p.text = "المساعد الذكي المتكامل لمتاجر سلة والتجارة الإلكترونية السعودية\nإعداد وتطوير: فريق إدارة وتطوير الأعمال"
    p.font.name = "Arial"
    p.font.size = Pt(12)
    p.font.color.rgb = COLOR_TEXT_MUTED
    p.alignment = PP_ALIGN.CENTER

    # ==========================================
    # SLIDE 2: EXECUTIVE SUMMARY & VISION
    # ==========================================
    slide2 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide2)
    add_header(slide2, "نظرة عامة • EXECUTIVE SUMMARY", "رؤية منصة سند وموقعها في السوق", "تمكين متاجر التجارة الإلكترونية السعودية من أتمتة خدمة العملاء والمبيعات بكفاءة تفوق 98%")

    cards_data_s2 = [
      ("🎯 الرؤية الاستراتيجية", "أن نكون المنصة السحابية رقم #1 في المملكة والخليج لأتمتة عمليات خدمة العملاء والمبيعات ومتابعة الشحن لمتاجر سلة وزد بالذكاء الاصطناعي.", COLOR_BRAND),
      ("💡 القيمة المضافة للتاجر", "توفير أكثر من 40 ساعة عمل أسبوعياً، وخفض تكاليف التوظيف بنسبة 70%، مع رفع سرعة الاستجابة للعملاء إلى أقل من 1.2 ثانية على مدار الساعة.", COLOR_GREEN),
      ("🚀 النمو المستهدف", "مضاعفة قاعدة المتاجر النشطة، وتوسيع الشراكات مع منصات التجارة الإلكترونية وشركات الشحن الكبرى لتعزيز الحصة السوقية.", COLOR_AMBER)
    ]

    for i, (title, desc, border_col) in enumerate(cards_data_s2):
        c = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8 + i*3.95), Inches(2.0), Inches(3.7), Inches(4.8))
        c.fill.solid()
        c.fill.fore_color.rgb = COLOR_SURFACE
        c.line.color.rgb = border_col
        c.line.width = Pt(1.5)
        
        tf = c.text_frame
        tf.word_wrap = True
        tf.margin_top = Inches(0.3)
        tf.margin_left = Inches(0.25)
        tf.margin_right = Inches(0.25)
        
        p = tf.paragraphs[0]
        p.text = title
        p.font.name = "Arial"
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = COLOR_TEXT_WHITE
        
        p_desc = tf.add_paragraph()
        p_desc.text = "\n" + desc
        p_desc.font.name = "Arial"
        p_desc.font.size = Pt(13)
        p_desc.font.color.rgb = COLOR_TEXT_MUTED

    # ==========================================
    # SLIDE 3: COMPREHENSIVE SWOT ANALYSIS (2x2 Matrix)
    # ==========================================
    slide3 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide3)
    add_header(slide3, "التحليل الاستراتيجي • SWOT MATRIX", "مصفوفة تحليل SWOT لمنصة سند", "تقييم دقيق لنقاط القوة والضعف الداخلية والفرص والتهديدات الخارجية")

    swot_items = [
      # (Col, Row, Title, Items, BorderColor, BgIcon)
      (0, 0, "💪 نقاط القوة (STRENGTHS)", [
        "تكامل مباشر مع سلة وشركات الشحن (SMSA, Aramex, OTO)",
        "دعم كامل وفائق للهجة والسياق المحلي السعودي",
        "قاعدة معرفة صارمة (Zero Hallucination) بدون هلوسة",
        "صندوق وارد موحد يدمج واتساب ويب وانستغرام وسلة",
        "تكلفة تشغيلية منخفضة وهوامش ربح سحابية عالية"
      ], COLOR_GREEN),
      
      (1, 0, "⚠️ نقاط الضعف (WEAKNESSES)", [
        "علامة تجارية حديثة العهد مقارنة بالشركات العالمية",
        "الاعتماد الحالي على WhatsApp Web قبل اعتماد Cloud API الرسمي",
        "الحاجة لإضافة تكاملات شحن إضافية (DHL, J&T)",
        "محدودية فريق المبيعات المباشر في المرحلة الحالية"
      ], COLOR_AMBER),

      (0, 1, "🚀 الفرص السوقية (OPPORTUNITIES)", [
        "سوق ضخم متسارع يضم أكثر من 50,000 متجر على سلة وزد",
        "ارتفاع تكلفة موظفي خدمة العملاء في السعودية (>4,000 ر.س)",
        "التوسع لدعم منصة زد (Zid) ومنصة شوبيفاي (Shopify)",
        "إطلاق ميزة استرداد السلات المتروكة لزيادة مبيعات التاجر",
        "بناء برنامج شركاء مع وكالات إدارة المتاجر والتسويق"
      ], COLOR_BRAND_LIGHT),

      (1, 1, "🛡️ التهديدات والتحديات (THREATS)", [
        "دخول منصات عالمية (Wati / ManyChat) للشرق الأوسط",
        "تغييرات محتملة في سياسات Meta و WhatsApp API",
        "حرب أسعار من أدوات الشات بوت التقليدية منخفضة الجودة",
        "تراجع وعي بعض التجار بقدرات الذكاء الاصطناعي الحديث"
      ], COLOR_RED)
    ]

    for col, row, title, items, color in swot_items:
        left = Inches(0.8 + col * 5.95)
        top = Inches(2.0 + row * 2.5)
        width = Inches(5.75)
        height = Inches(2.35)

        card = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR_SURFACE
        card.line.color.rgb = color
        card.line.width = Pt(1.5)

        tf = card.text_frame
        tf.word_wrap = True
        tf.margin_top = Inches(0.15)
        tf.margin_left = Inches(0.2)
        tf.margin_right = Inches(0.2)

        p = tf.paragraphs[0]
        p.text = title
        p.font.name = "Arial"
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = color

        for item in items:
            p_item = tf.add_paragraph()
            p_item.text = "• " + item
            p_item.font.name = "Arial"
            p_item.font.size = Pt(10.5)
            p_item.font.color.rgb = COLOR_TEXT_WHITE

    # ==========================================
    # SLIDE 4: COMPETITOR ANALYSIS & MATRIX
    # ==========================================
    slide4 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide4)
    add_header(slide4, "المشهد التنافسي • COMPETITIVE LANDSCAPE", "تحليل ومقارنة المنافسين في السوق", "مقارنة سند مع الحلول العالمية والتقليدية والموظف البشري")

    # Table of Competitors
    rows, cols = 5, 5
    table_shape = slide4.shapes.add_table(rows, cols, Inches(0.8), Inches(2.0), Inches(11.7), Inches(4.8))
    table = table_shape.table

    # Column widths
    table.columns[0].width = Inches(2.5)
    table.columns[1].width = Inches(2.3)
    table.columns[2].width = Inches(2.3)
    table.columns[3].width = Inches(2.3)
    table.columns[4].width = Inches(2.3)

    table_data = [
      ["معيار المقارنة", "⚡ سند (Sanad AI)", "🌐 شات بوت عالمي (Wati/ManyChat)", "🏢 أنظمة تقليدية (Zendesk/Intercom)", "👤 الموظف البشري التقليدي"],
      ["الربط المباشر بسلة والشحن", "✅ مدمج ومباشر ولحظي", "❌ يتطلب وسيط مثل Zapier", "❌ معقد ويحتاج مبرمجين", "⚠️ يدوي وبطيء مع أخطاء"],
      ["فهم اللهجة والسياق السعودي", "✅ ذكاء لغوي متقدم 100%", "⚠️ روبوت بقوالب جامدة", "⚠️ ترجمة آلية ضعيفة", "✅ فهم ممتاز لكن محدود بالوقت"],
      ["سرعة الاستجابة والتغطية", "⚡ 1.2 ثانية • على مدار 24/7", "⏱️ فوري لكن إجابات محدودة", "⏱️ ساعات في الدعم العادي", "⏳ 15-45 دقيقة (8 ساعات فقط)"],
      ["التكلفة الشهرية", "💰 تبدأ من 149 ر.س فقط", "💳 تبدأ من 400+ ر.س (قوالب)", "💳 مكلفة جداً (>1,500 ر.س)", "💸 راتب 3,500 - 5,000 ر.س"]
    ]

    for r_idx, row in enumerate(table_data):
        for c_idx, cell_value in enumerate(row):
            cell = table.cell(r_idx, c_idx)
            cell.text = cell_value
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            p.font.name = "Arial"
            
            if r_idx == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLOR_BRAND if c_idx == 1 else COLOR_SURFACE_ALT
                p.font.bold = True
                p.font.size = Pt(11)
                p.font.color.rgb = COLOR_TEXT_WHITE
            else:
                cell.fill.solid()
                if c_idx == 1:
                    cell.fill.fore_color.rgb = RGBColor(30, 20, 55) # Highlight Sanad Column
                    p.font.bold = True
                    p.font.color.rgb = COLOR_BRAND_ACCENT
                else:
                    cell.fill.fore_color.rgb = COLOR_SURFACE
                    p.font.color.rgb = COLOR_TEXT_WHITE
                p.font.size = Pt(10.5)

    # ==========================================
    # SLIDE 5: COMPETITIVE ADVANTAGE & USP
    # ==========================================
    slide5 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide5)
    add_header(slide5, "الميزة التنافسية • COMPETITIVE MOAT", "ما الذي يجعل (سند) الخيار الأول لمتجر سلة؟", "3 ركائز أساسية تمنح سند ميزة تنافسية مستدامة يصعب تقليدها")

    usps = [
      ("🇸🇦 التخصيص الكامل للسوق السعودي", "بناء النموذج اللغوي ليفهم مصطلحات التجارة السعودية (شحن سمسا، تقسيط تابي وتمارا، استبدال، استرجاع، مقاسات العبايات والأزياء)، مما يرفع ثقة العميل بنسبة 99%.", COLOR_BRAND),
      ("⚡ التفعيل الفوري (Plug & Play)", "بدون أي خبرة برمجية أو إعدادات معقدة؛ التاجر يربط متجره على سلة وواتساب في أقل من 5 دقائق ويبدأ المساعد في العمل فوراً وبشكل ذاتي.", COLOR_GREEN),
      ("📈 عائد استثماري مباشر (High ROI)", "سند لا يكتفي بالرد على الشكاوى؛ بل يسترد السلات المتروكة، ويقترح المنتجات البديلة، ويوفر أكثر من 3,500 ريال شهرياً على التاجر من أول أسبوع.", COLOR_AMBER)
    ]

    for i, (title, desc, color) in enumerate(usps):
        c = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8 + i*3.95), Inches(2.0), Inches(3.7), Inches(4.8))
        c.fill.solid()
        c.fill.fore_color.rgb = COLOR_SURFACE
        c.line.color.rgb = color
        c.line.width = Pt(1.5)

        tf = c.text_frame
        tf.word_wrap = True
        tf.margin_top = Inches(0.3)
        tf.margin_left = Inches(0.25)
        tf.margin_right = Inches(0.25)

        p = tf.paragraphs[0]
        p.text = title
        p.font.name = "Arial"
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = COLOR_TEXT_WHITE

        p_desc = tf.add_paragraph()
        p_desc.text = "\n" + desc
        p_desc.font.name = "Arial"
        p_desc.font.size = Pt(12)
        p_desc.font.color.rgb = COLOR_TEXT_MUTED

    # ==========================================
    # SLIDE 6: NEXT QUARTER OBJECTIVES (OKRs)
    # ==========================================
    slide6 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide6)
    add_header(slide6, "الأهداف الاستراتيجية • NEXT QUARTER OKRs", "أهداف ومؤشرات أداء الكوارتر القادم", "خارطة طريق واضحة وقابلة للقياس لتحقيق قفزة نوعية في النمو والإيرادات")

    okrs = [
      ("📈 الهدف 1: النمو واكتساب المتاجر", [
        "الوصول إلى 350+ متجر سلة نشط ومشترك",
        "تحقيق نسبة نمو شهري في الاشتراكات (MoM) لا تقل عن 25%",
        "زيادة معدل تحويل التجربة المجانية إلى اشتراك مدفوع إلى 35%"
      ], COLOR_BRAND),

      ("💰 الهدف 2: الأداء المالي والإيرادات", [
        "تحقيق 120,000 ريال إيراد شهري متكرر (MRR)",
        "رفع متوسط قيمة الاشتراك للمتجر عبر باقة Pro (349 ر.س)",
        "الحفاظ على معدل الإلغاء (Churn Rate) أقل من 3% شهرياً"
      ], COLOR_GREEN),

      ("⚙️ الهدف 3: تطوير المنتج والذكاء الاصطناعي", [
        "خفض زمن الاستجابة إلى أقل من 0.9 ثانية لكل رسالة",
        "رفع دقة استعلامات الشحن والمخزون إلى 99.7%",
        "إطلاق محرك استرداد السلات المتروكة الذكي بنسبة استرداد >18%"
      ], COLOR_AMBER),

      ("🤝 الهدف 4: التوسع والشراكات", [
        "إطلاق الربط الرسمي مع منصة زد (Zid Integration)",
        "تفعيل الربط المباشر مع WhatsApp Cloud API الرسمي",
        "بناء 10 شراكات استراتيجية مع وكالات إدارة المتاجر الإلكترونية"
      ], COLOR_BLUE)
    ]

    for idx, (title, points, color) in enumerate(okrs):
        col = idx % 2
        row = idx // 2
        left = Inches(0.8 + col * 5.95)
        top = Inches(2.0 + row * 2.5)

        card = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(5.75), Inches(2.35))
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR_SURFACE
        card.line.color.rgb = color
        card.line.width = Pt(1.5)

        tf = card.text_frame
        tf.word_wrap = True
        tf.margin_top = Inches(0.15)
        tf.margin_left = Inches(0.2)
        tf.margin_right = Inches(0.2)

        p = tf.paragraphs[0]
        p.text = title
        p.font.name = "Arial"
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = color

        for pt in points:
            p_pt = tf.add_paragraph()
            p_pt.text = "✔ " + pt
            p_pt.font.name = "Arial"
            p_pt.font.size = Pt(11)
            p_pt.font.color.rgb = COLOR_TEXT_WHITE

    # ==========================================
    # SLIDE 7: EXECUTION TIMELINE & ROADMAP
    # ==========================================
    slide7 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide7)
    add_header(slide7, "خارطة طريق التنفيذ • ACTION PLAN", "الجدول الزمني ومراحل التنفيذ خلال الكوارتر", "توزيع المهام والمستهدفات على مدار 3 أشهر لضمان التنفيذ الدقيق")

    months = [
      ("الشهر الأول (Month 1)\nالتركيز على الانتشار والتحسين", [
        "إطلاق حملة تسويقية تستهدف متاجر سلة الأكثر مبيعاً",
        "تحسين واجهة المستخدم وتجربة التسجيل في 1-Click",
        "رفع كفاءة تتبع شحنات سمسا وأرامكس وأوتو",
        "الهدف: الوصول إلى 150 متجر نشط"
      ], COLOR_BRAND),

      ("الشهر الثاني (Month 2)\nالتوسع في الميزات والمنصات", [
        "إطلاق تكامل منصة زد (Zid) رسمياً",
        "تفعيل ميزة استرداد السلات المتروكة للواتساب",
        "إطلاق حملات إعادة الاستهداف وعروض الدفع السنوي",
        "الهدف: الوصول إلى 250 متجر نشط"
      ], COLOR_GREEN),

      ("الشهر الثالث (Month 3)\nالشراكات والترقية المؤسسية", [
        "تدشين برنامج الشركاء والمسوقين بالعمولة",
        "إطلاق الباقة المؤسسية (Enterprise) للبراندات الكبرى",
        "مراجعة شاملة لنسب الرضا ومعدلات الإلغاء",
        "الهدف: تجاوز 350 متجر و 120k MRR"
      ], COLOR_AMBER)
    ]

    for i, (m_title, m_points, color) in enumerate(months):
        c = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8 + i*3.95), Inches(2.0), Inches(3.7), Inches(4.8))
        c.fill.solid()
        c.fill.fore_color.rgb = COLOR_SURFACE
        c.line.color.rgb = color
        c.line.width = Pt(1.5)

        tf = c.text_frame
        tf.word_wrap = True
        tf.margin_top = Inches(0.3)
        tf.margin_left = Inches(0.2)
        tf.margin_right = Inches(0.2)

        p = tf.paragraphs[0]
        p.text = m_title
        p.font.name = "Arial"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = color

        for pt in m_points:
            p_pt = tf.add_paragraph()
            p_pt.text = "\n• " + pt
            p_pt.font.name = "Arial"
            p_pt.font.size = Pt(11)
            p_pt.font.color.rgb = COLOR_TEXT_WHITE

    # ==========================================
    # SLIDE 8: SUMMARY & NEXT STEPS
    # ==========================================
    slide8 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide8)
    add_header(slide8, "الخاتمة والتوصيات • SUMMARY & NEXT STEPS", "أهم التوصيات لتسريع النمو وتحقيق الأهداف", "خطة العمل الفورية لفرق التسويق، المنتج، والمبيعات")

    summary_cards = [
      ("🎯 تسويق موجه ومكثف", "التركيز على إبراز التوفير المالي (ROI) وسرعة استجابة سند في المحتوى التسويقي لمجتمعات تجار سلة وزد.", COLOR_BRAND),
      ("💎 تجربة عميل استثنائية", "المتابعة الاستباقية مع كل متجر خلال فترة الـ 14 يوماً المجانية لضمان الربط الناجح وتحقيق أعلى نسبة تحويل للاشتراك.", COLOR_GREEN),
      ("🚀 التطوير المستمر للذكاء الاصطناعي", "تغذية النماذج ببيانات متجددة من قطاع التجارة الإلكترونية لضمان بقاء سند الأذكى والأكثر دقة في السوق السعودي.", COLOR_AMBER)
    ]

    for i, (title, desc, color) in enumerate(summary_cards):
        c = slide8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8 + i*3.95), Inches(2.0), Inches(3.7), Inches(3.6))
        c.fill.solid()
        c.fill.fore_color.rgb = COLOR_SURFACE
        c.line.color.rgb = color
        c.line.width = Pt(1.5)

        tf = c.text_frame
        tf.word_wrap = True
        tf.margin_top = Inches(0.25)
        tf.margin_left = Inches(0.2)
        tf.margin_right = Inches(0.2)

        p = tf.paragraphs[0]
        p.text = title
        p.font.name = "Arial"
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = COLOR_TEXT_WHITE

        p_desc = tf.add_paragraph()
        p_desc.text = "\n" + desc
        p_desc.font.name = "Arial"
        p_desc.font.size = Pt(12)
        p_desc.font.color.rgb = COLOR_TEXT_MUTED

    # Bottom Banner
    banner = slide8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(5.9), Inches(11.7), Inches(1.0))
    banner.fill.solid()
    banner.fill.fore_color.rgb = COLOR_BRAND
    banner.line.fill.background()
    tf = banner.text_frame
    p = tf.paragraphs[0]
    p.text = "✨ سند (Sanad AI) — شريكك الذكي للريادة في التجارة الإلكترونية السعودية ✨"
    p.font.name = "Arial"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLOR_TEXT_WHITE
    p.alignment = PP_ALIGN.CENTER

    # Save presentation
    prs.save(output_path)
    print(f"Presentation saved successfully at: {output_path}")

if __name__ == "__main__":
    output_file = r"c:\Users\Dell\Desktop\Support-Hub-AI\Sanad_Strategic_Plan_Q_Next.pptx"
    create_deck(output_file)
